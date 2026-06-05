from pptx import Presentation
prs = Presentation(r"D:\cupt作品赛\答辩PPT_环形喷泉_附录版.pptx")
slide = prs.slides[7]  # Slide 8 (0-indexed)
print("=== Slide 8 ===")
for shape in slide.shapes:
    if shape.has_text_frame:
        t = shape.text_frame.text.strip()
        if t:
            print(f"\nName='{shape.name}', top={shape.top}, h={shape.height}, w={shape.width}, left={shape.left}")
            print(f"  Text: '{t[:80]}'")
