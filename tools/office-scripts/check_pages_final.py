from pptx import Presentation
prs = Presentation(r"D:\cupt作品赛\答辩PPT_环形喷泉_附录版.pptx")

with open(r"C:\Users\lenovo\check_pages_final.txt", "w", encoding="utf-8") as f:
    for i in range(5, 12):
        slide = prs.slides[i]
        f.write(f"\nSlide {i+1}:")
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t and len(t) <= 3 and t.replace(' ','').replace('-','').strip():
                    # Short text that could be a page number
                    is_num = t.isdigit()
                    if is_num or (len(t) <= 3):
                        f.write(f"\n  [{shape.name}] '{t}' (is_digit={t.isdigit()})")

    f.write("\n\n--- Checking slides 23-28 ---\n")
    for i in range(22, 28):
        slide = prs.slides[i]
        f.write(f"\nSlide {i+1}:")
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t.isdigit() and len(t) <= 2:
                    f.write(f" [{shape.name}]='{t}'")

print("Done")
