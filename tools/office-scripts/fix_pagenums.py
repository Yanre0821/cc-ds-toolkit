"""Fix page numbers: slides 8-27 (old 7-26) need +1 offset."""
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

prs = Presentation(r"D:\cupt作品赛\答辩PPT_环形喷泉_附录版.pptx")

GRAY = RGBColor(0x80, 0x80, 0x80)
updated = 0

for i in range(7, 27):  # Slides 8-27 (0-indexed: 7-26)
    slide = prs.slides[i]
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t.isdigit() and len(t) <= 2:
                old = int(t)
                # Old slide 7 (now at position 8) had page num 7, should be 8
                # Old slide 8 (now at position 9) had page num 8, should be 9
                # etc.
                expected_old = i  # old page num (before insert)
                if old == expected_old:
                    new = old + 1
                    # Update the text
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.text.strip() == str(old):
                                run.text = str(new)
                                updated += 1
                                break

print(f"Updated {updated} page numbers")
prs.save(r"D:\cupt作品赛\答辩PPT_环形喷泉_附录版.pptx")
print("Saved")
