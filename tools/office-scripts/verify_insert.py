from pptx import Presentation

prs = Presentation(r"D:\cupt作品赛\答辩PPT_环形喷泉_附录版.pptx")

with open(r"C:\Users\lenovo\verify_insert.txt", "w", encoding="utf-8") as f:
    f.write(f"Total slides: {len(prs.slides)}\n")

    for i in [5, 6, 7, 8]:
        slide = prs.slides[i]
        f.write(f"\n{'='*60}\nSLIDE {i+1}\n{'='*60}\n")
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    f.write(f"\n[{shape.name}] top={shape.top}\n{t[:200]}\n")

    # Also check appendix divider
    for i in [27, 28]:
        slide = prs.slides[i]
        f.write(f"\n{'='*60}\nSLIDE {i+1}\n{'='*60}\n")
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    f.write(f"\n[{shape.name}]\n{t[:100]}\n")

print("Done")
